"""Deterministic real-format fixtures for complex Office dogfood."""
from __future__ import annotations

import hashlib
import re
import struct
from io import BytesIO
from pathlib import Path
from typing import Protocol, cast
from zipfile import ZIP_DEFLATED, ZIP_STORED, ZipFile, ZipInfo


class _Chart(Protocol):
    def add_data(self, data: object, *, titles_from_data: bool = False) -> None: ...


class _TextShape(Protocol):
    text: str
    width: int
    height: int


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def artifact(path: Path) -> dict[str, str]:
    digest = sha256(path)
    return {"artifact_id": digest, "content_hash": digest, "media_type": "application/octet-stream", "uri": str(path.resolve()), "sensitivity": "internal", "acl_fingerprint": "complex-dogfood"}


def part_hashes(path: Path) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {name: hashlib.sha256(archive.read(name)).hexdigest() for name in archive.namelist()}


def _rewrite(path: Path, changes: dict[str, bytes]) -> None:
    with ZipFile(path) as source:
        parts = {name: source.read(name) for name in source.namelist()}
    parts.update(changes)
    with ZipFile(path, "w") as output:
        for name, payload in parts.items():
            info = ZipInfo(name, (1980, 1, 1, 0, 0, 0))
            info.compress_type = ZIP_STORED if name == "mimetype" else ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            output.writestr(info, payload)


def docx(path: Path) -> Path:
    from docx import Document
    from docx.enum.text import WD_BREAK

    document = Document()
    for page in range(1, 4):
        _ = document.add_heading(f"복합 품질 보고서 / Complex QA {page}", level=1)
        _ = document.add_paragraph("한국어와 English 본문, comments, fields, revisions evidence.")
        if page == 1:
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text, table.cell(0, 1).text = "항목", "Value"
            table.cell(1, 0).text, table.cell(1, 1).text = "고객", "서울"
            _ = document.add_paragraph("PLACEHOLDER")
        document.sections[-1].header.paragraphs[0].text = "BIRKIN / 비르킨"
        document.sections[-1].footer.paragraphs[0].text = "Internal QA footer"
        if page < 3:
            document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    document.save(str(path))
    with ZipFile(path) as package:
        body = package.read("word/document.xml")
        types = package.read("[Content_Types].xml")
        rels = package.read("word/_rels/document.xml.rels")
    match = re.search(rb"<w:p[^>]*>(?:(?!</w:p>).)*PLACEHOLDER.*?</w:p>", body, re.DOTALL)
    if match is None:
        raise AssertionError("DOCX placeholder missing")
    rich = (b'<w:p><w:sdt><w:sdtPr><w:tag w:val="customer"/></w:sdtPr><w:sdtContent>'
            b'<w:r><w:t>PLACE</w:t></w:r><w:r><w:t>HOLDER</w:t></w:r>'
            b'</w:sdtContent></w:sdt><w:commentRangeStart w:id="0"/>'
            b'<w:ins w:id="7" w:author="QA"><w:r><w:t> inserted</w:t></w:r></w:ins>'
            b'<w:del w:id="8" w:author="QA"><w:r><w:delText> deleted</w:delText></w:r></w:del>'
            b'<w:fldSimple w:instr="DATE"><w:r><w:t>2026-08-16</w:t></w:r></w:fldSimple>'
            b'<w:commentRangeEnd w:id="0"/><w:r><w:commentReference w:id="0"/></w:r></w:p>')
    comments = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                '<w:comment w:id="0" w:author="QA"><w:p><w:r><w:t>검토 comment</w:t></w:r></w:p></w:comment></w:comments>').encode()
    types = types.replace(b"</Types>", b'<Override PartName="/word/comments.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml"/></Types>')
    rels = rels.replace(b"</Relationships>", b'<Relationship Id="rIdDogfoodComments" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" Target="comments.xml"/></Relationships>')
    _rewrite(path, {"word/document.xml": body.replace(match.group(0), rich, 1), "word/comments.xml": comments, "word/_rels/document.xml.rels": rels, "[Content_Types].xml": types, "customXml/untouched.bin": b"DOCX-SENTINEL"})
    return path


def xlsx(path: Path) -> Path:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.chart._chart import ChartBase
    from openpyxl.styles import Font, PatternFill
    from openpyxl.worksheet.worksheet import Worksheet

    book = Workbook()
    data = cast(Worksheet, book.active)
    data.title = "Data"
    data.append(["지역 / Region", "매출 / Revenue"])
    for row in [("서울", 10), ("Busan", 20), ("제주", 30)]:
        data.append(row)
    data["A1"].font = Font(bold=True, color="FFFFFF")
    data["A1"].fill = PatternFill("solid", fgColor="1F4E78")
    data.row_dimensions[3].hidden = True
    data.column_dimensions["D"].hidden = True
    chart = cast(_Chart, BarChart())
    chart.add_data(Reference(data, min_col=2, min_row=1, max_row=4), titles_from_data=True)
    data.add_chart(cast(ChartBase, chart), "F2")
    hidden = cast(Worksheet, book.create_sheet("Hidden"))
    hidden.sheet_state = "hidden"
    hidden["A1"] = "hidden evidence"
    secret = cast(Worksheet, book.create_sheet("VeryHidden"))
    secret.sheet_state = "veryHidden"
    secret["A1"] = "very hidden evidence"
    book.save(path)
    with ZipFile(path) as package:
        sheet = package.read("xl/worksheets/sheet1.xml")
    marker = b"</sheetData>"
    cells = (b'<row r="6"><c r="B6"><f>SUM(B2:B4)</f><v>999</v></c></row>'
             b'<row r="7"><c r="B7"><f t="array" ref="B7:B9">_xlfn.UNIQUE(B2:B4)</f><v>10</v></c></row>'
             b'<row r="8"><c r="B8"><f>\'[external.xlsx]Data\'!B2</f><v>77</v></c></row>'
             b'<row r="9"><c r="B9" t="e"><f>1/0</f><v>#DIV/0!</v></c></row>')
    sheet = sheet.replace(b'<c r="B2" t="n">', b'<c r="B2">', 1)
    _rewrite(path, {"xl/worksheets/sheet1.xml": sheet.replace(marker, cells + marker), "custom/untouched.bin": b"XLSX-SENTINEL"})
    return path


def pptx(path: Path) -> Path:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Pt

    logo = path.with_suffix(".png")
    Image.new("RGB", (80, 40), "#1F4E78").save(logo)
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    title = slide.shapes.title
    if title is None:
        raise AssertionError("PPTX title placeholder missing")
    title.text = "BIRKIN / 복합 프레젠테이션"
    body = cast(_TextShape, cast(object, slide.placeholders[1]))
    body.text = "PLACEHOLDER"
    body.width, body.height = Inches(1), Inches(0.15)
    _ = Pt(42)
    _ = slide.shapes.add_picture(str(logo), Inches(8.7), Inches(6.7), width=Inches(2))
    notes = slide.notes_slide.notes_text_frame
    if notes is None:
        raise AssertionError("PPTX notes text frame missing")
    notes.text = "발표자 노트 / speaker notes"
    second = deck.slides.add_slide(deck.slide_layouts[5])
    second_title = second.shapes.title
    if second_title is None:
        raise AssertionError("PPTX second title placeholder missing")
    second_title.text = "Evidence"
    deck.save(str(path))
    logo.unlink()
    with ZipFile(path) as package:
        xml = package.read("ppt/slides/slide1.xml")
    xml = xml.replace(b"<a:t>PLACEHOLDER</a:t>", b'<a:t>PLACE</a:t></a:r><a:r><a:rPr><a:latin typeface="Definitely Missing QA Font"/></a:rPr><a:t>HOLDER</a:t>', 1)
    _rewrite(path, {"ppt/slides/slide1.xml": xml, "custom/untouched.bin": b"PPTX-SENTINEL"})
    return path


def hwpx(path: Path) -> Path:
    section0 = ('<hs:sec xmlns:hs="hs" xmlns:hp="hp"><hp:p id="P0"><hp:field id="customer">'
                '<hp:run><hp:t>PLACE</hp:t></hp:run><hp:run><hp:t>HOLDER</hp:t></hp:run>'
                '</hp:field><hp:t> 한국어 English</hp:t></hp:p><hp:tbl id="T1"><hp:tr><hp:tc><hp:p><hp:t>표</hp:t></hp:p></hp:tc></hp:tr></hp:tbl></hs:sec>').encode()
    parts = {"mimetype": b"application/hwp+zip", "META-INF/manifest.xml": b"<manifest><item>template</item></manifest>", "Contents/content.hpf": b"<opf><section id='0'/><section id='1'/></opf>", "Contents/header.xml": b"<head><styles>Branded</styles></head>", "Contents/section0.xml": section0, "Contents/section1.xml": '<hs:sec xmlns:hs="hs" xmlns:hp="hp"><hp:p id="P1"><hp:t>두 번째 section</hp:t></hp:p></hs:sec>'.encode(), "Contents/untouched.xml": b"<sentinel>HWPX</sentinel>"}
    with ZipFile(path, "w") as package:
        for name, payload in parts.items():
            info = ZipInfo(name, (1980, 1, 1, 0, 0, 0))
            info.compress_type = ZIP_STORED if name == "mimetype" else ZIP_DEFLATED
            package.writestr(info, payload)
    return path


def _pdf(objects: list[bytes]) -> bytes:
    data = bytearray(b"%PDF-1.7\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, body in enumerate(objects, 1):
        offsets.append(len(data))
        data.extend(f"{number} 0 obj\n".encode() + body + b"\nendobj\n")
    xref = len(data)
    data.extend(f"xref\n0 {len(objects)+1}\n0000000000 65535 f \n".encode())
    for offset in offsets[1:]:
        data.extend(f"{offset:010d} 00000 n \n".encode())
    data.extend(f"trailer\n<< /Size {len(objects)+1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode())
    return bytes(data)


def pdfs(folder: Path) -> dict[str, Path]:
    content = b"BT /F1 12 Tf 72 720 Td (Complex native PDF) Tj ET"
    base = [b"<< /Type /Catalog /Pages 2 0 R >>", b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>", b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>", b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>", b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream"]
    variants = {"native": _pdf(base)}
    image = b"\0\0\0"
    drawing = b"q 100 0 0 100 0 0 cm /Im0 Do Q"
    variants["scanned"] = _pdf([base[0], base[1], b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 100 100] /Resources << /XObject << /Im0 4 0 R >> >> /Contents 5 0 R >>", b"<< /Type /XObject /Subtype /Image /Width 1 /Height 1 /ColorSpace /DeviceRGB /BitsPerComponent 8 /Length 3 >>\nstream\n" + image + b"\nendstream", b"<< /Length " + str(len(drawing)).encode() + b" >>\nstream\n" + drawing + b"\nendstream"])
    variants["form"] = _pdf([b"<< /Type /Catalog /Pages 2 0 R /AcroForm 6 0 R >>", base[1], base[2][:-3] + b" /Annots [7 0 R] >>", base[3], base[4], b"<< /Fields [7 0 R] >>", b"<< /Type /Annot /Subtype /Widget /FT /Tx /T (name) /Rect [0 0 50 20] /P 3 0 R >>"])
    variants["active"] = _pdf([b"<< /Type /Catalog /Pages 2 0 R /OpenAction 6 0 R /Names << /JavaScript << /Names [(startup) 7 0 R] >> >> >>", *base[1:], b"<< /Type /Action /S /Launch /F (never.bin) >>", b"<< /Type /Action /S /JavaScript /JS (app.alert('never')) >>"])
    variants["signed"] = _pdf([b"<< /Type /Catalog /Pages 2 0 R /AcroForm 6 0 R /Perms << /DocMDP 7 0 R >> >>", *base[1:], b"<< /Fields [8 0 R] /SigFlags 3 >>", b"<< /Type /Sig /Filter /Adobe.PPKLite /SubFilter /adbe.pkcs7.detached /ByteRange [0 0 0 0] /Contents <01020304> >>", b"<< /FT /Sig /T (Approval) /V 7 0 R >>"])
    paths: dict[str, Path] = {}
    for name, payload in variants.items():
        paths[name] = folder / f"{name}.pdf"
        _ = paths[name].write_bytes(payload)
    from pypdf import PdfReader, PdfWriter
    writer = PdfWriter(clone_from=PdfReader(BytesIO(variants["native"]), strict=True))
    writer.encrypt("correct horse", algorithm="RC4-128")
    paths["encrypted"] = folder / "encrypted.pdf"
    with paths["encrypted"].open("wb") as output:
        _ = writer.write(output)
    return paths


def odf_and_legacy(folder: Path) -> dict[str, Path]:
    result: dict[str, Path] = {}
    for ext, mime, root in (("odt", "application/vnd.oasis.opendocument.text", "text"), ("ods", "application/vnd.oasis.opendocument.spreadsheet", "spreadsheet"), ("odp", "application/vnd.oasis.opendocument.presentation", "presentation")):
        path = folder / f"identity.{ext}"
        with ZipFile(path, "w") as package:
            package.writestr("mimetype", mime)
            package.writestr("content.xml", f'<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"><office:body><office:{root}/></office:body></office:document-content>')
        result[ext] = path
    free, end, fat_marker = 0xFFFFFFFF, 0xFFFFFFFE, 0xFFFFFFFD
    for ext, identity in (("doc", "WordDocument"), ("xls", "Workbook"), ("ppt", "PowerPoint Document"), ("hwp", "HWP Document File")):
        encoded = (identity + "\0").encode("utf-16le")
        entry = bytearray(128)
        entry[:len(encoded)] = encoded
        struct.pack_into("<H", entry, 64, len(encoded))
        entry[66] = 2
        struct.pack_into("<III", entry, 68, free, free, free)
        struct.pack_into("<I", entry, 116, end)
        root = bytearray(128)
        name = "Root Entry\0".encode("utf-16le")
        root[:len(name)] = name
        struct.pack_into("<H", root, 64, len(name))
        root[66] = 5
        struct.pack_into("<III", root, 68, free, free, free)
        struct.pack_into("<I", root, 116, end)
        header = bytearray(512)
        header[:8] = bytes.fromhex("d0cf11e0a1b11ae1")
        struct.pack_into("<HHHH", header, 24, 0x3E, 3, 0xFFFE, 9)
        struct.pack_into("<H", header, 32, 6)
        struct.pack_into("<IIIIIIIII", header, 40, 0, 1, 0, 0, 4096, end, 0, end, 0)
        struct.pack_into("<109I", header, 76, 1, *([free] * 108))
        fat = [end, fat_marker] + [free] * 126
        path = folder / f"identity.{ext}"
        _ = path.write_bytes(bytes(header) + (bytes(root) + bytes(entry)).ljust(512, b"\0") + struct.pack("<128I", *fat))
        result[ext] = path
    result["rtf"] = folder / "identity.rtf"
    _ = result["rtf"].write_bytes(b"{\\rtf1\\ansi\\ansicpg949 Complex \\field Korean-English}")
    return result
