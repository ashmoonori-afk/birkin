"""Literal offline complex-document dogfood with bounded JSON evidence."""
from __future__ import annotations

import argparse
import importlib
import json
import os
import shutil
import sys
from collections.abc import Iterator
from pathlib import Path
from typing import Protocol, cast

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

_pdf_adapter = importlib.import_module("birkin.office.adapters.pdf")
_pptx_adapter = importlib.import_module("birkin.office.adapters.pptx")
_xlsx_adapter = importlib.import_module("birkin.office.adapters.xlsx")
_approvals = importlib.import_module("birkin.approvals")
_conversion_audit = importlib.import_module("birkin.office.conversion_audit")
_errors = importlib.import_module("birkin.office.errors")
_legacy_conversion = importlib.import_module("birkin.office.legacy_conversion")
_legacy_preflight = importlib.import_module("birkin.office.legacy_preflight")
_legacy_types = importlib.import_module("birkin.office.legacy_types")
_service = importlib.import_module("birkin.office.service")
_tools = importlib.import_module("birkin.tools")
_tool_types = importlib.import_module("birkin.tools._types")
_complex_fixtures = importlib.import_module("script.qa.office_complex_fixtures")
_complex_report = importlib.import_module("script.qa.office_complex_report")

PdfAdapter = _pdf_adapter.PdfAdapter
PptxAdapter = _pptx_adapter.PptxAdapter
XlsxAdapter = _xlsx_adapter.XlsxAdapter
approve = _approvals.approve
LOSS_CATEGORIES = _conversion_audit.LOSS_CATEGORIES
DocumentError = _errors.DocumentError
DocumentService = _service.DocumentService
convert_legacy = _legacy_conversion.convert_legacy
preflight_legacy = _legacy_preflight.preflight_legacy
LegacyConversionRequest = _legacy_types.LegacyConversionRequest
LegacyEnginePin = _legacy_types.LegacyEnginePin
LegacyRefusal = _legacy_types.LegacyRefusal
build_registry = _tools.build_registry
ToolContext = _tool_types.ToolContext
artifact = _complex_fixtures.artifact
docx = _complex_fixtures.docx
hwpx = _complex_fixtures.hwpx
odf_and_legacy = _complex_fixtures.odf_and_legacy
part_hashes = _complex_fixtures.part_hashes
pdfs = _complex_fixtures.pdfs
pptx = _complex_fixtures.pptx
sha256 = _complex_fixtures.sha256
xlsx = _complex_fixtures.xlsx
write_bounded_report = _complex_report.write_bounded_report

FORMATS = ("docx", "xlsx", "pptx", "pdf", "hwpx")
CREATE: dict[str, dict[str, object]] = {
    "docx": {"paragraphs": ["Created DOCX / 생성 문서"]},
    "xlsx": {"sheets": [{"name": "Created", "rows": [["값", 1]]}]},
    "pptx": {"slides": [{"title": "Created PPTX", "body": "생성됨"}]},
    "pdf": {"paragraphs": ["Created PDF / 생성 문서"]},
}
class _IterText(Protocol):
    def itertext(self) -> Iterator[str]: ...


class _TextShape(Protocol):
    text: str


REQUESTS = {
    "docx": "Update this Word document",
    "xlsx": "Update this Excel spreadsheet",
    "pptx": "Update this PowerPoint presentation",
    "pdf": "Update this PDF document",
    "hwpx": "Update this HWPX document",
}
PATCH: dict[str, dict[str, object]] = {
    "docx": {
        "locator": {"format": "docx", "index": 2},
        "value": "수정된 고객 / Modified customer",
    },
    "xlsx": {"cell": "B2", "value": 42},
    "pptx": {"placeholder_idx": 1, "value": "수정됨 / Modified"},
    "hwpx": {"field": "customer", "value": "수정된 고객 / Modified customer"},
}


def _error(error: DocumentError) -> dict[str, object]:
    return cast(dict[str, object], error.envelope()["error"])


def _reopen_validation(format_name: str, path: Path) -> tuple[dict[str, object], bool]:
    if format_name == "docx":
        from docx import Document

        document = Document(str(path))
        matched = "Modified customer" in "".join(
            cast(_IterText, document.element).itertext()
        )
        return {"status": "ok", "parser": "python-docx"}, matched
    if format_name == "xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(path, data_only=False, read_only=False)
        try:
            matched = cast(object, workbook["Data"]["B2"].value) == 42
        finally:
            workbook.close()
        return {"status": "ok", "parser": "openpyxl"}, matched
    if format_name == "pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        text = "\n".join(
            cast(_TextShape, cast(object, shape)).text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        return {"status": "ok", "parser": "python-pptx"}, "Modified" in text
    if format_name == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(path, strict=True)
        return {"status": "ok", "parser": "pypdf", "pages": len(reader.pages)}, len(reader.pages) == 1
    if format_name == "hwpx":
        from zipfile import ZipFile

        from defusedxml import ElementTree

        with ZipFile(path) as package:
            root = ElementTree.fromstring(package.read("Contents/section0.xml"))
        matched = "Modified customer" in "".join(root.itertext())
        return {"status": "ok", "parser": "defusedxml"}, matched
    raise AssertionError(f"no reopen parser registered for {format_name}")


def _compact_hash_evidence(before: dict[str, str], after: dict[str, str]) -> dict[str, object]:
    common = before.keys() & after.keys()
    changed = sorted(name for name in common if before[name] != after[name])
    protected = sorted(name for name in common if name.endswith(".rels") or "vba" in name.lower() or "signature" in name.lower() or name.startswith("custom/"))
    return {
        "part_count_before": len(before), "part_count_after": len(after),
        "changed_parts": changed, "untouched_parts_preserved": all(before[name] == after[name] for name in common - set(changed)),
        "relationships_macros_signatures": {name: {"before": before[name], "after": after[name], "preserved": before[name] == after[name]} for name in protected},
    }


def run(output_dir: Path) -> dict[str, object]:
    jail = output_dir.expanduser().resolve()
    if jail.exists() and any(jail.iterdir()):
        raise ValueError("output directory must be absent or empty")
    jail.mkdir(parents=True, exist_ok=True)
    temporary = jail / ".fixture-staging"
    temporary.mkdir()
    office_home = jail / "office"
    sources = office_home / "sources"
    sources.mkdir(parents=True)
    os.environ["BIRKIN_HOME"] = str(jail)
    registry = build_registry(ToolContext(cfg={"spill_threshold": 1_000_000}, client=None, cwd=jail), include={"documents"})
    service = DocumentService(office_home)
    coordinated = office_home / "coordinated"
    coordinated.mkdir()
    tools_exercised: set[str] = set()

    def execute(name: str, data: dict[str, object]) -> tuple[dict[str, object], bool]:
        tools_exercised.add(name)
        result = registry.execute(name, data)
        if not isinstance(result.content, str):
            raise TypeError(f"{name} returned non-JSON content")
        return cast(dict[str, object], json.loads(result.content)), result.is_error

    def call(name: str, data: dict[str, object], *, refused: bool = False) -> dict[str, object]:
        body, is_error = execute(name, data)
        if is_error != refused:
            raise AssertionError(f"{name} unexpected result: {body}")
        return body

    def coordinate(format_name: str, source: dict[str, str], operation: dict[str, object]) -> dict[str, str]:
        destination = coordinated / f"modified-{format_name}.{format_name}"
        proposed = call(
            "office_job_request",
            {
                "request": REQUESTS[format_name],
                "source": source,
                "outcome": f"Apply the approved {format_name} edit",
                "operations": [operation],
                "destination": str(destination),
            },
        )
        approved = approve(
            cast(str, proposed["id"]),
            approved_by="system:qa",
            approved_via="qa:script",
        )
        if approved["ok"] is not True:
            raise AssertionError(f"office_job_request approval failed: {approved}")
        return artifact(destination)

    paths = {
        "docx": docx(sources / "complex.docx"),
        "xlsx": xlsx(sources / "complex.xlsx"),
        "pptx": pptx(sources / "complex.pptx"),
        "hwpx": hwpx(sources / "complex.hwpx"),
    }
    pdf_paths = pdfs(sources)
    paths["pdf"] = pdf_paths["native"]
    identities = odf_and_legacy(sources)
    records: dict[str, dict[str, object]] = {}
    expected_refusals: list[dict[str, object]] = []

    inventory = cast(list[dict[str, object]], call("list_document_adapters", {})["adapters"])
    adapter_by_format = {str(item["format"]): item for item in inventory}
    template_plan = service.fill_template(
        artifact(paths["hwpx"]),
        [{"key": "customer", "value": "템플릿 고객 / Template"}],
        fields=[{"key": "customer", "kind": "field", "field": "customer"}],
        output_name="planned.hwpx",
    )

    for fmt in FORMATS:
        source_path = paths[fmt]
        source = artifact(source_path)
        source_digest = sha256(source_path)
        before_parts = {} if fmt == "pdf" else part_hashes(source_path)
        content = CREATE.get(fmt, {"bindings": {"customer": "생성된 고객 / Created"}})
        create_data: dict[str, object] = {"format": fmt, "content": content, "output_name": f"created-{fmt}.{fmt}"}
        if fmt == "hwpx":
            create_data["template"] = source
        try:
            create_result = service.create_document(**create_data)
            created = cast(dict[str, str], create_result["draft_artifact"])
            creation: dict[str, object] = {"status": "ok", "sha256": created["content_hash"]}
        except DocumentError as error:
            create_error = _error(error)
            creation = {"status": "unavailable", "error": create_error}
            expected_refusals.append({"format": fmt, "operation": "create", **create_error})
        inspection = call("inspect_document", {"source": source})
        extraction = call("extract_document", {"source": source})
        validation = call("validate_artifact", {"artifact": source})
        comparison = call("compare_documents", {"left": source, "right": source})
        preview = call("render_artifact", {"artifact": source, "output_format": "structured_preview"})
        visual = call("render_artifact", {"artifact": source, "output_format": "png"}, refused=True)
        visual_error = cast(dict[str, object], visual["error"])
        expected_refusals.append({"format": fmt, "operation": "visual_render", **visual_error})
        budget = {category: 100 for category in LOSS_CATEGORIES}
        converted = cast(dict[str, str], service.convert_document(source, target_format="txt", output_name=f"complex-{fmt}.txt", loss_budget=budget)["draft_artifact"])
        primary = source
        mutation: dict[str, object]
        if fmt == "pdf":
            try:
                _ = service.apply_document_patch(
                    source,
                    {"operations": [{"type": "body_edit"}]},
                    expected_source_sha256=source_digest,
                    output_name="modified.pdf",
                    dry_run=False,
                )
                raise AssertionError("PDF body edit unexpectedly succeeded")
            except DocumentError as error:
                refused = _error(error)
            mutation = {"status": "expected_refusal", "error": refused}
            expected_refusals.append({"format": "pdf", "operation": "body_edit", **refused})
            preservation: dict[str, object] = {"source_immutable": sha256(source_path) == source_digest, "signature_parts": "not_applicable_native_fixture"}
        else:
            if fmt == "docx":
                modified = coordinate(fmt, source, PATCH[fmt])
            else:
                patched = service.apply_document_patch(
                    source,
                    {"operations": [PATCH[fmt]]},
                    expected_source_sha256=source_digest,
                    output_name=f"modified-{fmt}.{fmt}",
                    dry_run=False,
                )
                modified = cast(dict[str, str], patched["draft_artifact"])
            primary = modified
            changed = call("compare_documents", {"left": source, "right": modified})
            if changed["equal"] is True:
                raise AssertionError(f"{fmt} mutation did not change the artifact")
            preservation = _compact_hash_evidence(before_parts, part_hashes(Path(modified["uri"])))
            mutation = {"status": "ok", "output_sha256": modified["content_hash"]}
        if sha256(source_path) != source_digest:
            raise AssertionError(f"{fmt} source changed")
        reopened, expected_content_match = _reopen_validation(
            fmt, Path(primary["uri"])
        )
        if not expected_content_match:
            raise AssertionError(f"{fmt} reopened but edited content did not match")
        structure = cast(dict[str, object], inspection["structure"])["inventory"]
        records[fmt] = {
            "source_sha256_before": source_digest, "source_sha256_after": sha256(source_path),
            "capabilities": adapter_by_format[fmt]["capabilities"], "modified_artifact": primary["uri"],
            "operations": {"create": creation, "inspect": "ok", "extract": {"status": "ok", "spans": len(cast(list[object], extraction["spans"]))}, "modify": mutation, "validate": {"status": "ok", "checks": len(cast(list[object], validation["checks"]))}, "diff": {"status": "ok", "self_equal": comparison["equal"]}, "structured_preview": {"status": "preview", "visual_proof": preview["visual_proof"]}, "convert_txt": {"status": "ok", "sha256": converted["content_hash"]}, "visual_render": {"status": "unavailable", "code": visual_error["code"]}},
            "structure": structure, "preservation": preservation,
            "reopen_validation": reopened,
            "expected_content_match": expected_content_match,
        }

    xlsx_audit = XlsxAdapter().audit_formulas(paths["xlsx"])
    recalc = XlsxAdapter().recalculate(paths["xlsx"])
    layout = PptxAdapter().audit_layout(paths["pptx"])
    records["hwpx"]["template_plan"] = template_plan
    records["xlsx"]["formula_evidence"] = {"count": len(xlsx_audit["cells"]), "states": sorted({cell["cache_status"] for cell in xlsx_audit["cells"]}), "dynamic": xlsx_audit["dynamic_arrays"], "external": xlsx_audit["external_links"], "recalculation": recalc}
    records["pptx"]["layout_evidence"] = {"warnings": layout["warnings"], "fonts": layout["fonts"], "media": layout["media"], "visual_verification": {"status": "not-run", "adapter_evidence": layout["visual_verification"]}}

    pdf_states: dict[str, object] = {}
    for name, path in pdf_paths.items():
        registered = call("inspect_document", {"source": artifact(path)})
        structure = cast(dict[str, object], registered["structure"])
        state = cast(dict[str, object], structure["inventory"])
        pdf_states[name] = {key: state.get(key) for key in ("content_type", "form_type", "encrypted", "signed", "active_content", "signatures")}
        cast(dict[str, object], pdf_states[name])["credential_required"] = state.get("credential_required", state.get("password_required"))
    pdf_refusal_cases = [("scanned", "extract", lambda: PdfAdapter().extract(pdf_paths["scanned"])), ("form", "form_fill", lambda: PdfAdapter().fill(pdf_paths["form"], {"name": "Ada"})), ("active", "mutation", lambda: PdfAdapter().patch(pdf_paths["active"], {"type": "body_edit"})), ("signed", "mutation", lambda: PdfAdapter().patch(pdf_paths["signed"], {"type": "body_edit"})), ("encrypted", "extract", lambda: PdfAdapter().extract(pdf_paths["encrypted"]))]
    typed_pdf_refusals: list[dict[str, object]] = []
    for variant, operation, action in pdf_refusal_cases:
        try:
            _ = action()
            raise AssertionError(f"PDF {variant} {operation} unexpectedly succeeded")
        except DocumentError as error:
            typed_pdf_refusals.append({"variant": variant, "operation": operation, **_error(error)})
    records["pdf"]["state_cases"] = pdf_states
    records["pdf"]["inspection_surface"] = "registered_inspect_document"
    records["pdf"]["typed_refusals"] = typed_pdf_refusals
    records["pdf"]["not_run"] = {"ocr": {"status": "unavailable", "reason": "pdf_image_only_requires_ocr"}, "form_write": {"status": "unsupported"}, "visual_render": {"status": "unavailable"}}

    unsupported: dict[str, object] = {}
    for ext in ("odt", "ods", "odp", "hwp"):
        refusal = call("inspect_document", {"source": artifact(identities[ext])}, refused=True)
        error = cast(dict[str, object], refusal["error"])
        unsupported[ext] = {"path": str(identities[ext]), "identity_sha256": sha256(identities[ext]), "status": "unsupported", "error": error}
    try:
        _ = preflight_legacy(identities["hwp"])
        raise AssertionError("binary HWP unexpectedly accepted as a legacy Office input")
    except LegacyRefusal as refusal:
        hwp_evidence = cast(dict[str, object], unsupported["hwp"])
        hwp_evidence["strict_identity_refusal"] = refusal.receipt.to_dict()
    legacy: dict[str, object] = {}
    filters = {"doc": ("MS Word 97", "Office Open XML Text", "docx"), "xls": ("MS Excel 97", "Calc MS Excel 2007 XML", "xlsx"), "ppt": ("MS PowerPoint 97", "Impress MS PowerPoint 2007 XML", "pptx"), "rtf": ("Rich Text Format", "Office Open XML Text", "docx")}
    for ext, (input_filter, output_filter, target) in filters.items():
        source = identities[ext]
        preflight = preflight_legacy(source)
        request = LegacyConversionRequest(target_format=target, engine=LegacyEnginePin("libreoffice", "24.2.7.2", input_filter, output_filter))
        receipt = convert_legacy(source, temporary / f"legacy-{ext}.{target}", request)
        if (
            receipt.status != "refused"
            or receipt.reason_code != "external_engine_forbidden"
        ):
            raise AssertionError(
                f"legacy {ext} external engine was not permanently refused"
            )
        legacy[ext] = {"identity": preflight.to_dict(), "conversion": receipt.to_dict(), "source_immutable": sha256(source) == preflight.source_sha256}
    shutil.rmtree(temporary)
    cleanup = {"temporary_paths_removed": [str(temporary)], "temporary_paths_remaining": [str(temporary)] if temporary.exists() else [], "artifacts_retained": True}
    ok = all(item["source_sha256_before"] == item["source_sha256_after"] for item in records.values()) and not cleanup["temporary_paths_remaining"] and all(cast(dict[str, object], item)["status"] == "unsupported" for item in unsupported.values())
    return {"ok": ok, "jail": str(jail), "formats": records, "unsupported_identities": unsupported, "legacy": legacy, "expected_refusals": expected_refusals, "cleanup_receipt": cleanup, "tools_exercised": sorted(tools_exercised)}


def main() -> int:
    parser = argparse.ArgumentParser()
    _ = parser.add_argument("--output-dir", required=True, type=Path)
    args = parser.parse_args()
    report: dict[str, object]
    try:
        report = run(cast(Path, args.output_dir))
    except (AssertionError, DocumentError, ImportError, LookupError, OSError, RuntimeError, TypeError, ValueError) as error:
        report = {"ok": False, "error": {"type": type(error).__name__, "message": str(error)}}
    if report["ok"] is True:
        report = write_bounded_report(report, cast(Path, args.output_dir))
    print(json.dumps(report, ensure_ascii=False, sort_keys=True, default=str))
    return 0 if report["ok"] is True else 1


if __name__ == "__main__":
    raise SystemExit(main())
