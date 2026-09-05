from __future__ import annotations

import hashlib
import importlib
import zipfile
from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from birkin.office.errors import DocumentError, DocumentErrorCode
from birkin.office.service import DocumentService
from tests.office.fixture_builders import build_hwpx_template

BUSINESS_TEMPLATES = {
    "weekly_report": {
        "title": "주간 업무 보고",
        "period": "2026-09-01 ~ 2026-09-05",
        "summary": "출시 준비 완료",
        "achievements": ["품질 검사 완료"],
        "metrics": [["지표", "값"], ["완료", "3"]],
    },
    "meeting_notes": {
        "title": "제품 회의록",
        "date": "2026-09-05",
        "summary": "출시 범위를 확정함",
        "decisions": ["DOCX부터 출시"],
        "actions": [["담당", "할 일"], ["민수", "검증"]],
    },
    "work_proposal": {
        "title": "문서 자동화 제안",
        "problem": "반복 작성 시간이 큼",
        "proposal": "검증된 양식을 사용",
        "benefits": ["작성 시간 단축"],
        "costs": [["항목", "비용"], ["개발", "2일"]],
    },
}


def _hash(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _artifact(path: Path) -> dict[str, str]:
    digest = _hash(path)
    return {
        "artifact_id": digest,
        "content_hash": digest,
        "media_type": "application/octet-stream",
        "uri": str(path),
        "sensitivity": "internal",
        "acl_fingerprint": "a" * 64,
    }


@pytest.mark.parametrize(
    ("fmt", "plan"),
    [
        ("docx", {"paragraphs": ["Contract", "Approved"]}),
        ("xlsx", {"sheets": [{"name": "Data", "rows": [["A", 1]]}]}),
        ("pptx", {"slides": [{"title": "Contract", "body": "Approved"}]}),
        ("pdf", {"paragraphs": ["Contract", "Approved"]}),
    ],
)
def test_typed_blank_plan_emits_reopenable_file_and_complete_receipt(
    tmp_path: Path, fmt: str, plan: dict[str, object]
) -> None:
    result = DocumentService(tmp_path).create_document(
        format=fmt, content=plan, output_name=f"contract.{fmt}"
    )
    output = Path(result["draft_artifact"]["uri"])
    assert result["format"] == fmt
    assert result["creation_mode"] == "blank_authoring"
    assert result["source_sha256"] is None
    assert result["template_sha256"] is None
    assert result["output_sha256"] == _hash(output)
    assert result["receipt"]["operation"] == "document_create"
    assert all(item["passed"] for item in result["validation_evidence"])
    assert result["capability_limits"] and result["fidelity_limits"]
    if fmt == "docx":
        assert Document(str(output)).paragraphs[1].text == "Approved"
    elif fmt == "xlsx":
        worksheet = load_workbook(output, read_only=True)["Data"]
        assert worksheet.cell(1, 2).value == 1
    elif fmt == "pptx":
        title = Presentation(str(output)).slides[0].shapes.title
        assert title is not None and title.text == "Contract"
    else:
        assert "Approved" in (PdfReader(output).pages[0].extract_text() or "")


@pytest.mark.parametrize(
    ("fmt", "plan"),
    [
        ("docx", {"paragraphs": ["ok"], "unknown": True}),
        ("xlsx", {"sheets": [{"name": "bad/name", "rows": [[1]]}]}),
        ("pptx", {"slides": [{"title": 3}]}),
        ("pdf", {"paragraphs": []}),
    ],
)
def test_invalid_plan_never_creates_destination(
    tmp_path: Path, fmt: str, plan: dict[str, object]
) -> None:
    service = DocumentService(tmp_path)
    with pytest.raises(DocumentError) as caught:
        _ = service.create_document(format=fmt, content=plan, output_name=f"bad.{fmt}")
    assert caught.value.code is DocumentErrorCode.INVALID_INPUT
    assert not (tmp_path / "artifacts" / "drafts" / f"bad.{fmt}").exists()


def test_missing_backend_is_typed_and_plan_validation_precedes_import(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    def missing(_name: str) -> None:
        raise ImportError("not installed")

    monkeypatch.setattr(importlib, "import_module", missing)
    service = DocumentService(tmp_path)
    with pytest.raises(DocumentError) as invalid:
        _ = service.create_document(
            format="docx",
            content={"paragraphs": ["valid"], "extra": True},
            output_name="invalid.docx",
        )
    assert invalid.value.code is DocumentErrorCode.INVALID_INPUT
    with pytest.raises(DocumentError) as unavailable:
        _ = service.create_document(
            format="docx",
            content={"paragraphs": ["valid"]},
            output_name="missing.docx",
        )
    assert unavailable.value.code is DocumentErrorCode.CAPABILITY_UNAVAILABLE
    assert "install_hint" in unavailable.value.details
    assert not (tmp_path / "artifacts" / "drafts" / "missing.docx").exists()


def test_output_guards_reject_escape_existing_and_unsupported(tmp_path: Path) -> None:
    service = DocumentService(tmp_path)
    cases = [
        ("docx", "../escape.docx", DocumentErrorCode.INVALID_INPUT),
        ("odt", "unsupported.odt", DocumentErrorCode.UNSUPPORTED_FORMAT),
    ]
    for fmt, name, code in cases:
        with pytest.raises(DocumentError) as caught:
            _ = service.create_document(
                format=fmt, content={"paragraphs": ["x"]}, output_name=name
            )
        assert caught.value.code is code
    existing = tmp_path / "artifacts" / "drafts" / "existing.docx"
    _ = existing.write_bytes(b"keep")
    with pytest.raises(DocumentError) as caught:
        _ = service.create_document(
            format="docx",
            content={"paragraphs": ["x"]},
            output_name="existing.docx",
        )
    assert caught.value.code is DocumentErrorCode.OUTPUT_EXISTS
    assert existing.read_bytes() == b"keep"


def test_stale_template_identity_and_same_destination_are_rejected(tmp_path: Path) -> None:
    service = DocumentService(tmp_path)
    template_path = build_hwpx_template(tmp_path / "stale-hash.hwpx")
    stale = _artifact(template_path)
    _ = template_path.write_bytes(template_path.read_bytes() + b"changed")
    with pytest.raises(DocumentError) as changed:
        _ = service.create_document(
            format="hwpx",
            content={"bindings": {"customer": "Ada"}},
            output_name="stale-hash-out.hwpx",
            template=stale,
        )
    assert changed.value.code is DocumentErrorCode.SOURCE_CHANGED
    same = build_hwpx_template(
        tmp_path / "artifacts" / "drafts" / "same.hwpx"
    )
    with pytest.raises(DocumentError) as collision:
        _ = service.create_document(
            format="hwpx",
            content={"bindings": {"customer": "Ada"}},
            output_name="same.hwpx",
            template=_artifact(same),
        )
    assert collision.value.code is DocumentErrorCode.OUTPUT_EXISTS


def test_hwpx_derivation_is_hash_bound_preserves_parts_and_reports_mode(
    tmp_path: Path,
) -> None:
    template_path = build_hwpx_template(tmp_path / "trusted.hwpx")
    template = _artifact(template_path)
    before = template_path.read_bytes()
    result = DocumentService(tmp_path).create_document(
        format="hwpx",
        content={
            "bindings": {
                "customer": {"value": "Ada", "expected_text": "PLACEHOLDER"}
            }
        },
        output_name="derived.hwpx",
        template=template,
    )
    output = Path(result["draft_artifact"]["uri"])
    assert template_path.read_bytes() == before
    assert result["creation_mode"] == "template_derivation"
    assert result["source_sha256"] == template["content_hash"]
    assert result["template_sha256"] == template["content_hash"]
    with zipfile.ZipFile(output) as archive:
        assert b"Ada" in archive.read("Contents/section0.xml")
        assert archive.read("Contents/opaque.xml") == b'<x:unknown xmlns:x="urn:opaque" a=" 1 "/>'


@pytest.mark.parametrize(
    ("mutator", "code"),
    [
        ("ambiguous", DocumentErrorCode.AMBIGUOUS_LOCATOR),
        ("stale", DocumentErrorCode.PRECONDITION_FAILED),
        ("active", DocumentErrorCode.POLICY_DENIED),
        ("signature", DocumentErrorCode.POLICY_DENIED),
    ],
)
def test_unsafe_or_unbound_template_is_refused_before_output(
    tmp_path: Path, mutator: str, code: DocumentErrorCode
) -> None:
    template_path = build_hwpx_template(tmp_path / f"{mutator}.hwpx")
    if mutator != "stale":
        with zipfile.ZipFile(template_path, "a", zipfile.ZIP_STORED) as archive:
            if mutator == "ambiguous":
                section = archive.read("Contents/section0.xml")
                archive.writestr("Contents/section1.xml", section)
            elif mutator == "active":
                archive.writestr("Contents/Scripts/default.js", b"never execute")
            else:
                archive.writestr("META-INF/signatures.xml", b"<signatures/>")
    content: dict[str, object] = {
        "bindings": {
            "customer": {"value": "Ada", "expected_text": "STALE"}
        }
    }
    with pytest.raises(DocumentError) as caught:
        _ = DocumentService(tmp_path).create_document(
            format="hwpx",
            content=content,
            output_name=f"{mutator}-out.hwpx",
            template=_artifact(template_path),
        )
    assert caught.value.code is code
    assert not (tmp_path / "artifacts" / "drafts" / f"{mutator}-out.hwpx").exists()


def test_template_risk_requires_specific_consent_and_is_warned(tmp_path: Path) -> None:
    template_path = build_hwpx_template(tmp_path / "active.hwpx")
    with zipfile.ZipFile(template_path, "a", zipfile.ZIP_STORED) as archive:
        archive.writestr("Contents/Scripts/default.js", b"never execute")
    result = DocumentService(tmp_path).create_document(
        format="hwpx",
        content={
            "bindings": {"customer": "Ada"},
            "allow_active_content": True,
        },
        output_name="consented.hwpx",
        template=_artifact(template_path),
    )
    assert any("never executed" in warning for warning in result["warnings"])


@pytest.mark.parametrize(("name", "values"), BUSINESS_TEMPLATES.items())
def test_approved_business_templates_create_structured_reopenable_docx(
    tmp_path: Path, name: str, values: dict[str, object]
) -> None:
    result = DocumentService(tmp_path).create_document(
        format="docx",
        content={"business_template": {
            "name": name,
            "version": "1.0",
            "values": values,
            "sources": {"title": "user://request"},
        }},
        output_name=f"{name}.docx",
    )
    reopened = Document(result["draft_artifact"]["uri"])
    assert reopened.paragraphs[0].text == values["title"]
    assert reopened.tables[0].cell(1, 1).text
    assert result["business_template"] == {
        "name": name,
        "version": "1.0",
        "profile_sha256": result["business_template"]["profile_sha256"],
        "required_fields": list({
            "weekly_report": ("title", "period", "summary"),
            "meeting_notes": ("title", "date", "summary"),
            "work_proposal": ("title", "problem", "proposal"),
        }[name]),
        "missing_fields": [],
        "unreplaced_fields": [],
        "sources": {"title": "user://request"},
        "layout_verified": False,
    }


@pytest.mark.parametrize(("name", "values"), BUSINESS_TEMPLATES.items())
def test_business_template_required_and_unreplaced_fields_fail_closed(
    tmp_path: Path, name: str, values: dict[str, object]
) -> None:
    incomplete = dict(values)
    del incomplete["title"]
    with pytest.raises(DocumentError, match="missing required values"):
        DocumentService(tmp_path).create_document(
            format="docx",
            content={"business_template": {"name": name, "version": "1.0", "values": incomplete}},
            output_name=f"missing-{name}.docx",
        )
    fields = tuple(values) + ("unbound",)
    template_path = build_hwpx_template(tmp_path / f"{name}.hwpx", fields)
    with pytest.raises(DocumentError, match="unreplaced fields"):
        DocumentService(tmp_path).create_document(
            format="hwpx",
            content={"business_template": {"name": name, "version": "1.0", "values": values}},
            output_name=f"unreplaced-{name}.hwpx",
            template=_artifact(template_path),
        )


@pytest.mark.parametrize(("name", "values"), BUSINESS_TEMPLATES.items())
def test_business_template_fills_reopenable_hwpx_and_records_template_hash(
    tmp_path: Path, name: str, values: dict[str, object]
) -> None:
    template_path = build_hwpx_template(tmp_path / f"source-{name}.hwpx", tuple(values))
    result = DocumentService(tmp_path).create_document(
        format="hwpx",
        content={"business_template": {"name": name, "version": "1.0", "values": values}},
        output_name=f"filled-{name}.hwpx",
        template=_artifact(template_path),
    )
    output = Path(result["draft_artifact"]["uri"])
    with zipfile.ZipFile(output) as archive:
        assert b"PLACEHOLDER" not in archive.read("Contents/section0.xml")
    assert result["template_sha256"] == _hash(template_path)
    assert result["business_template"]["layout_verified"] is False
