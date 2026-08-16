from __future__ import annotations

import zipfile
from pathlib import Path
from typing import cast

from docx import Document
from openpyxl import Workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.shapes.autoshape import Shape


def build_docx(path: Path, chunks: tuple[str, ...]) -> Path:
    document = Document()
    paragraph = document.add_paragraph()
    for chunk in chunks:
        _ = paragraph.add_run(chunk)
    document.save(str(path))
    return path


def build_xlsx(path: Path, text: str) -> Path:
    """Create a real workbook, then exercise OOXML shared-string storage."""
    workbook = Workbook()
    sheet = cast("Worksheet", workbook.active)
    sheet.title = "한글Data1"
    sheet["A1"] = text
    sheet["B1"] = '=CONCAT("원본", "-1")'
    sheet["C1"] = 7
    workbook.save(path)

    with zipfile.ZipFile(path) as source:
        parts = {name: source.read(name) for name in source.namelist()}
    sheet_xml = parts["xl/worksheets/sheet1.xml"]
    start = sheet_xml.index(b'<c r="A1"')
    end = sheet_xml.index(b"</c>", start) + 4
    parts["xl/worksheets/sheet1.xml"] = (
        sheet_xml[:start] + b'<c r="A1" t="s"><v>0</v></c>' + sheet_xml[end:]
    )
    parts["xl/sharedStrings.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        + b'<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        + b'count="1" uniqueCount="1"><si><r><t>'
        + _xml(text[: len(text) // 2])
        + b"</t></r><r><t>"
        + _xml(text[len(text) // 2 :])
        + b"</t></r></si></sst>"
    )
    rels = parts["xl/_rels/workbook.xml.rels"]
    parts["xl/_rels/workbook.xml.rels"] = rels.replace(
        b"</Relationships>",
        b'<Relationship Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings" '
        + b'Target="sharedStrings.xml" Id="rIdKorean"/></Relationships>',
    )
    types = parts["[Content_Types].xml"]
    parts["[Content_Types].xml"] = types.replace(
        b"</Types>",
        b'<Override PartName="/xl/sharedStrings.xml" '
        + b'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sharedStrings+xml"/></Types>',
    )
    _zip(path, parts)
    return path


def build_pptx(path: Path, chunks: tuple[str, ...]) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    subtitle = cast("Shape", slide.placeholders[1])
    if not subtitle.has_text_frame:
        raise TypeError("fixture subtitle must have a text frame")
    paragraph = subtitle.text_frame.paragraphs[0]
    _ = paragraph.clear()
    for chunk in chunks:
        paragraph.add_run().text = chunk
    presentation.save(str(path))
    return path


def build_hwpx(path: Path, chunks: tuple[str, ...]) -> Path:
    texts = b"".join(b"<hp:t>" + _xml(chunk) + b"</hp:t>" for chunk in chunks)
    section = (
        b'<?xml version="1.0" encoding="UTF-8"?><hs:sec '
        + b'xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section" '
        + b'xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">'
        + b'<hp:p id="korean-p"><hp:run>' + texts + b"</hp:run></hp:p>"
        + b'<hp:p id="field-p"><hp:run><hp:ctrl><hp:fieldBegin id="field-ko" name="'
        + "고객".encode()
        + b'" type="CLICK_HERE"/></hp:ctrl><hp:t>'
        + "필드값".encode()
        + b'</hp:t><hp:ctrl><hp:fieldEnd beginIDRef="field-ko"/></hp:ctrl>'
        + b"</hp:run></hp:p></hs:sec>"
    )
    parts = {
        "mimetype": b"application/hwp+zip",
        "META-INF/manifest.xml": b"<manifest/>",
        "Contents/section0.xml": section,
    }
    _zip(path, parts, first_stored=True)
    return path


def build_pdf(path: Path, text: str) -> Path:
    characters = {character: index for index, character in enumerate(dict.fromkeys(text), 1)}
    encoded = "".join(f"{characters[character]:04X}" for character in text)
    mappings = "\n".join(
        f'<{index:04X}> <{character.encode("utf-16-be").hex().upper()}>'
        for character, index in characters.items()
    )
    cmap = (
        "/CIDInit /ProcSet findresource begin 12 dict begin begincmap "
        "/CIDSystemInfo << /Registry (Adobe) /Ordering (UCS) /Supplement 0 >> def "
        "/CMapName /Korean def /CMapType 2 def 1 begincodespacerange <0000> <FFFF> "
        f"endcodespacerange {len(characters)} beginbfchar\n{mappings}\nendbfchar "
        "endcmap CMapName currentdict /CMap defineresource pop end end"
    ).encode()
    stream = f"BT /F1 6 Tf 36 720 Td <{encoded}> Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 7 0 R >>",
        b"<< /Type /Font /Subtype /Type0 /BaseFont /KoreanTest /Encoding /Identity-H /DescendantFonts [5 0 R] /ToUnicode 6 0 R >>",
        b"<< /Type /Font /Subtype /CIDFontType2 /BaseFont /KoreanTest /CIDSystemInfo << /Registry (Adobe) /Ordering (Identity) /Supplement 0 >> /DW 1000 /CIDToGIDMap /Identity >>",
        _stream(cmap),
        _stream(stream),
    ]
    payload = bytearray(b"%PDF-1.7\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, obj in enumerate(objects, 1):
        offsets.append(len(payload))
        payload += f"{index} 0 obj\n".encode() + obj + b"\nendobj\n"
    xref = len(payload)
    payload += f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode()
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode())
    payload += f"trailer << /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    _ = path.write_bytes(payload)
    return path


def _stream(data: bytes) -> bytes:
    return f"<< /Length {len(data)} >>\nstream\n".encode() + data + b"\nendstream"


def _xml(text: str) -> bytes:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").encode()


def _zip(path: Path, parts: dict[str, bytes], *, first_stored: bool = False) -> None:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for index, (name, payload) in enumerate(parts.items()):
            archive.writestr(name, payload, zipfile.ZIP_STORED if first_stored and index == 0 else zipfile.ZIP_DEFLATED)
