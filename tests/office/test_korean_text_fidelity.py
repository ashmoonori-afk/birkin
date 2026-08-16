from __future__ import annotations

import hashlib
import unicodedata
import zipfile
from collections.abc import Callable
from pathlib import Path
from typing import Protocol, cast

import pypdfium2 as pdfium  # pyright: ignore[reportMissingTypeStubs]
import pytest
from defusedxml import ElementTree
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.shapes.autoshape import Shape

from birkin.office.adapters.docx import DocxAdapter
from birkin.office.adapters.hwpx import HwpxAdapter
from birkin.office.adapters.hwpx_types import ParagraphLocator
from birkin.office.adapters.pdf import PdfAdapter
from birkin.office.adapters.pptx import PptxAdapter
from birkin.office.adapters.pptx_types import ShapeLocator
from birkin.office.adapters.xlsx import XlsxAdapter
from birkin.office.errors import DocumentError, DocumentErrorCode
from birkin.office.extract_package import extract_package_text
from tests.office.korean_fixtures import (
    build_docx,
    build_hwpx,
    build_pdf,
    build_pptx,
    build_xlsx,
)

NFC = "한글 가"
NFD = unicodedata.normalize("NFD", "각")
CP949 = "씨피949 원문".encode("cp949").decode("cp949")
SOURCE = f"{NFC} | {NFD} | ㄱㅏ | 👩‍💻 | Korea한글-2026 | 《괄호》,!? | e\u0301 | {CP949}"
EDITED = "수정 됨 | 한 | ㅋㅋ | 가족👨‍👩‍👧 | A한9 | ‘인용’ | o\u0308"
CHUNKS = (SOURCE[:7], SOURCE[7:19], SOURCE[19:])


class _PdfTextPage(Protocol):
    def get_text_bounded(self) -> str: ...


class _PdfPage(Protocol):
    def get_textpage(self) -> _PdfTextPage: ...


def _digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _part_hashes(path: Path) -> dict[str, str]:
    with zipfile.ZipFile(path) as archive:
        return {name: hashlib.sha256(archive.read(name)).hexdigest() for name in archive.namelist()}


def _changed_parts(source: Path, output: Path) -> set[str]:
    before, after = _part_hashes(source), _part_hashes(output)
    assert before.keys() == after.keys()
    return {name for name in before if before[name] != after[name]}


def test_corpus_has_distinct_codepoint_edge_cases_and_no_normalization_policy() -> None:
    assert unicodedata.is_normalized("NFC", NFC)
    assert not unicodedata.is_normalized("NFC", NFD)
    assert "ㄱㅏ" in SOURCE and "👩‍💻" in SOURCE and "e\u0301" in SOURCE
    assert unicodedata.normalize("NFC", SOURCE) != SOURCE
    assert CP949 == "씨피949 원문"


def test_docx_fragmented_runs_exact_edit_reextract_and_reopen(tmp_path: Path) -> None:
    source = build_docx(tmp_path / "한국어.docx", CHUNKS)
    output = tmp_path / "수정.docx"
    adapter = DocxAdapter()
    original_hash = _digest(source)
    first = adapter.inspect(source)
    assert first == adapter.inspect(source)
    node = next(item for item in first["paragraphs"] if item["text"] == SOURCE)
    assert extract_package_text(source, "docx") == [SOURCE]

    receipt = adapter.patch_text(
        source, output, node["locator"], EDITED,
        expected_text=SOURCE, expected_source_sha256=original_hash,
    )
    assert receipt["previous_text"] == SOURCE
    assert _digest(source) == original_hash
    assert _changed_parts(source, output) == {"word/document.xml"}
    assert extract_package_text(output, "docx") == [EDITED]
    reopened = Document(str(output))
    assert reopened.paragraphs[0].text == EDITED
    assert adapter.inspect(output)["paragraphs"][0]["text"] == EDITED


def test_xlsx_shared_string_survives_formula_edit_and_independent_reopen(tmp_path: Path) -> None:
    source = build_xlsx(tmp_path / "한국어.xlsx", SOURCE)
    output = tmp_path / "수정.xlsx"
    adapter = XlsxAdapter()
    original_hash = _digest(source)
    first = adapter.operation_inventory(source)
    assert first == adapter.operation_inventory(source)
    cell = next(
        item for item in first["cells"]
        if cast("dict[str, object]", item["locator"])["cell"] == "A1"
    )
    assert cell["storage"] == "shared_string" and cell["value"] == SOURCE
    assert extract_package_text(source, "xlsx")[0].split("\t")[0] == SOURCE

    receipt = adapter.patch_formula(
        source, output, {"sheet": "한글Data1", "cell": "B1"},
        'CONCAT("수정", "-2")', expected_formula='CONCAT("원본", "-1")',
        expected_source_sha256=original_hash,
    )
    assert receipt["cache_stale"] is True
    assert _digest(source) == original_hash
    assert _changed_parts(source, output) == {"xl/worksheets/sheet1.xml"}
    assert extract_package_text(output, "xlsx")[0].split("\t")[0] == SOURCE
    workbook = load_workbook(output, data_only=False)
    sheet = workbook["한글Data1"]
    assert cast("object", sheet["A1"].value) == SOURCE
    assert cast("object", sheet["B1"].value) == '=CONCAT("수정", "-2")'


def test_pptx_fragmented_placeholder_exact_edit_reextract_and_reopen(tmp_path: Path) -> None:
    source = build_pptx(tmp_path / "한국어.pptx", CHUNKS)
    output = tmp_path / "수정.pptx"
    adapter = PptxAdapter()
    original_hash = _digest(source)
    first = adapter.inventory(source)
    assert first == adapter.inventory(source)
    raw_locator = first["placeholders"][0]
    locator: ShapeLocator = {
        "part_uri": raw_locator["part_uri"],
        "shape_id": cast("str", raw_locator["shape_id"]),
    }
    assert extract_package_text(source, "pptx")[0] == SOURCE

    receipt = adapter.patch_text(
        source, output, locator, EDITED,
        expected_text=SOURCE, expected_source_sha256=original_hash,
    )
    assert receipt["previous_text"] == SOURCE
    assert _digest(source) == original_hash
    assert _changed_parts(source, output) == {"ppt/slides/slide1.xml"}
    assert extract_package_text(output, "pptx")[0] == EDITED
    reopened = Presentation(str(output))
    subtitle = cast("Shape", reopened.slides[0].placeholders[1])
    assert subtitle.text == EDITED


def test_hwpx_fragmented_text_and_field_exact_edit_reextract_and_reopen(tmp_path: Path) -> None:
    source = build_hwpx(tmp_path / "한국어.hwpx", CHUNKS)
    output = tmp_path / "수정.hwpx"
    adapter = HwpxAdapter()
    original_hash = _digest(source)
    first = adapter.inspect(source)
    assert first == adapter.inspect(source)
    assert first["paragraph_details"][0]["text"] == SOURCE
    assert first["field_details"][0]["key"] == "고객"
    assert extract_package_text(source, "hwpx") == [SOURCE, "필드값"]

    receipt = adapter.patch_paragraph_text(
        source, output, ParagraphLocator("Contents/section0.xml", "korean-p"),
        EDITED, expected_text=SOURCE, expected_source_sha256=original_hash,
    )
    assert receipt["previous_text"] == SOURCE
    assert _digest(source) == original_hash
    assert _changed_parts(source, output) == {"Contents/section0.xml"}
    assert extract_package_text(output, "hwpx") == [EDITED, "필드값"]
    with zipfile.ZipFile(output) as archive:
        root = ElementTree.fromstring(archive.read("Contents/section0.xml"))
    paragraphs = ["".join(node.itertext()) for node in root if node.tag.endswith("}p")]
    assert paragraphs == [EDITED, "필드값"]


def test_pdf_exact_native_extraction_independent_parser_and_edit_refusal(tmp_path: Path) -> None:
    source = build_pdf(tmp_path / "한국어.pdf", SOURCE)
    adapter = PdfAdapter()
    original_hash = _digest(source)
    assert adapter.inspect(source)["source_sha256"] == original_hash
    assert [span["text"] for span in adapter.extract(source)] == [SOURCE]
    document = pdfium.PdfDocument(source)
    page = cast("_PdfPage", cast("object", document[0]))
    text_page = page.get_textpage()
    assert text_page.get_text_bounded() == SOURCE

    with pytest.raises(DocumentError) as caught:
        adapter.patch(source, {"type": "body_edit", "value": EDITED})
    assert caught.value.code is DocumentErrorCode.UNSUPPORTED_EDIT
    assert _digest(source) == original_hash
    assert [span["text"] for span in adapter.extract(source)] == [SOURCE]


@pytest.mark.parametrize("bad", ["불가\x01", "불가\ud800"])
def test_illegal_xml_text_is_typed_refusal_without_output(tmp_path: Path, bad: str) -> None:
    docx = build_docx(tmp_path / "source.docx", (SOURCE,))
    docx_node = DocxAdapter().inspect(docx)["paragraphs"][0]
    pptx = build_pptx(tmp_path / "source.pptx", (SOURCE,))
    raw_pptx_locator = PptxAdapter().inventory(pptx)["placeholders"][0]
    pptx_locator: ShapeLocator = {
        "part_uri": raw_pptx_locator["part_uri"],
        "shape_id": cast("str", raw_pptx_locator["shape_id"]),
    }
    hwpx = build_hwpx(tmp_path / "source.hwpx", (SOURCE,))
    xlsx = build_xlsx(tmp_path / "source.xlsx", SOURCE)
    cases: list[tuple[Path, Path, Callable[[Path], object]]] = [
        (docx, tmp_path / "bad.docx", lambda output: DocxAdapter().patch_text(docx, output, docx_node["locator"], bad)),
        (pptx, tmp_path / "bad.pptx", lambda output: PptxAdapter().patch_text(pptx, output, pptx_locator, bad)),
        (hwpx, tmp_path / "bad.hwpx", lambda output: HwpxAdapter().patch_paragraph_text(hwpx, output, ParagraphLocator("Contents/section0.xml", "korean-p"), bad)),
        (xlsx, tmp_path / "bad.xlsx", lambda output: XlsxAdapter().patch_formula(xlsx, output, {"sheet": "한글Data1", "cell": "B1"}, bad)),
    ]
    for source, output, operation in cases:
        before = _digest(source)
        with pytest.raises(DocumentError) as caught:
            _ = operation(output)
        assert caught.value.code is DocumentErrorCode.INVALID_INPUT
        assert not output.exists() and _digest(source) == before
