from __future__ import annotations

import hashlib
import zipfile
from pathlib import Path
from typing import cast

import pytest
from pptx import Presentation

from birkin.office.adapters.pptx import PptxAdapter
from birkin.office.errors import DocumentError, DocumentErrorCode

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PR = "http://schemas.openxmlformats.org/package/2006/relationships"
C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
STRICT_P = "http://purl.oclc.org/ooxml/presentationml/main"
STRICT_A = "http://purl.oclc.org/ooxml/drawingml/main"


def _rels(*items: str) -> bytes:
    return (f'<Relationships xmlns="{PR}">' + "".join(items) + "</Relationships>").encode()


def _zip(path: Path, parts: dict[str, bytes]) -> Path:
    with zipfile.ZipFile(path, "w") as archive:
        for name, value in parts.items():
            info = zipfile.ZipInfo(name, (2025, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            archive.writestr(info, value)
    return path


def _rel(rid: str, kind: str, target: str, mode: str = "") -> str:
    external = ' TargetMode="External"' if mode else ""
    return f'<Relationship Id="{rid}" Type="{R}/{kind}" Target="{target}"{external}/>'


def _slide(one: bool) -> bytes:
    extras = ""
    if one:
        extras = (
            '<p:sp><p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr/><p:nvPr><p:ph idx="7"/></p:nvPr></p:nvSpPr>'
            '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr b="1"/><a:t>BRAND</a:t></a:r>'
            '<a:r><a:rPr i="1"/><a:t> TITLE</a:t></a:r></a:p></p:txBody></p:sp>'
            '<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="3" name="Table"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>'
            '<a:graphic><a:graphicData><a:tbl><a:tr h="100"><a:tc><a:txBody><a:p><a:r><a:rPr lang="ko-KR"/><a:t>A1</a:t></a:r></a:p></a:txBody></a:tc>'
            '<a:tc><a:txBody><a:p><a:r><a:rPr b="1"/><a:t>B1</a:t></a:r></a:p></a:txBody></a:tc></a:tr></a:tbl></a:graphicData></a:graphic></p:graphicFrame>'
            '<p:pic><p:nvPicPr><p:cNvPr id="4" name="Image"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr><p:blipFill><a:blip r:embed="rIdImage"/></p:blipFill></p:pic>'
            f'<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="5" name="Chart"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr><a:graphic><a:graphicData><c:chart xmlns:c="{C}" r:id="rIdChart"/></a:graphicData></a:graphic></p:graphicFrame>'
            '<p:transition advClick="1"/>'
        )
    return f'<p:sld xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}"><p:cSld><p:spTree>{extras}</p:spTree></p:cSld></p:sld>'.encode()


def _package(
    path: Path, *, sections: bool = True, shared_image_owner: str | None = None
) -> Path:
    content_types = (
        b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        b'<Default Extension="xml" ContentType="application/xml"/><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        b'<Default Extension="png" ContentType="image/png"/>'
        b'<Default Extension="xlsx" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"/>'
        b'<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
        b'<Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        b'<Override PartName="/ppt/slides/slide2.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        b'<Override PartName="/ppt/charts/chart1.xml" ContentType="application/vnd.openxmlformats-officedocument.drawingml.chart+xml"/>'
        b'</Types>'
    )
    section = '<p:sectionLst><p:section name="Keep" firstSlideId="256"/></p:sectionLst>' if sections else ""
    entries = {
        "[Content_Types].xml": content_types,
        "_rels/.rels": _rels(_rel("rId1", "officeDocument", "ppt/presentation.xml")),
        "ppt/presentation.xml": f'<p:presentation xmlns:p="{P}" xmlns:r="{R}"><p:sldMasterIdLst/><p:sldIdLst><p:sldId id="256" r:id="rId1"/><p:sldId id="257" r:id="rId2"/></p:sldIdLst><p:sldSz cx="9144000" cy="6858000" type="screen4x3"/>{section}</p:presentation>'.encode(),
        "ppt/_rels/presentation.xml.rels": _rels(_rel("rId1", "slide", "slides/slide1.xml"), _rel("rId2", "slide", "slides/slide2.xml")),
        "ppt/slides/slide1.xml": _slide(True),
        "ppt/slides/slide2.xml": _slide(False),
        "ppt/slides/_rels/slide1.xml.rels": _rels(_rel("rIdLayout", "slideLayout", "../slideLayouts/slideLayout1.xml"), _rel("rIdImage", "image", "../media/image1.png"), _rel("rIdChart", "chart", "../charts/chart1.xml"), _rel("rIdNotes", "notesSlide", "../notesSlides/notesSlide1.xml"), _rel("rIdComments", "comments", "../comments/comment1.xml")),
        "ppt/slides/_rels/slide2.xml.rels": _rels(_rel("rIdLayout", "slideLayout", "../slideLayouts/slideLayout1.xml")),
        "ppt/slideLayouts/slideLayout1.xml": f'<p:sldLayout xmlns:p="{P}"/>'.encode(),
        "ppt/slideLayouts/_rels/slideLayout1.xml.rels": _rels(_rel("rIdMaster", "slideMaster", "../slideMasters/slideMaster1.xml")),
        "ppt/slideMasters/slideMaster1.xml": f'<p:sldMaster xmlns:p="{P}"/>'.encode(),
        "ppt/slideMasters/_rels/slideMaster1.xml.rels": _rels(_rel("rIdTheme", "theme", "../theme/theme1.xml")),
        "ppt/theme/theme1.xml": f'<a:theme xmlns:a="{A}" name="Brand"/>'.encode(),
        "ppt/notesSlides/notesSlide1.xml": f'<p:notes xmlns:p="{P}" xmlns:a="{A}"><p:cSld><p:spTree><p:sp><p:nvSpPr><p:cNvPr id="9" name="Notes"/><p:cNvSpPr/><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:txBody><a:p><a:r><a:rPr lang="en-US"/><a:t>Speaker note</a:t></a:r></a:p></p:txBody></p:sp></p:spTree></p:cSld></p:notes>'.encode(),
        "ppt/charts/chart1.xml": f'<c:chartSpace xmlns:c="{C}"><c:chart><c:plotArea><c:barChart><c:ser><c:tx><c:v>Series</c:v></c:tx></c:ser></c:barChart></c:plotArea></c:chart></c:chartSpace>'.encode(),
        "ppt/charts/_rels/chart1.xml.rels": _rels(_rel("rIdWorkbook", "package", "../embeddings/Microsoft_Excel_Worksheet1.xlsx")),
        "ppt/embeddings/Microsoft_Excel_Worksheet1.xlsx": b"EMBEDDED-WORKBOOK-SENTINEL",
        "ppt/media/image1.png": b"\x89PNG\r\n\x1a\nORIGINAL",
        "ppt/comments/comment1.xml": b"<comments>KEEP</comments>",
        "customXml/item1.xml": b"<unknown>KEEP</unknown>",
    }
    if shared_image_owner is not None:
        owners = {
            "layout": ("ppt/slideLayouts/slideLayout1.xml", "ppt/slideLayouts/_rels/slideLayout1.xml.rels", "sldLayout"),
            "master": ("ppt/slideMasters/slideMaster1.xml", "ppt/slideMasters/_rels/slideMaster1.xml.rels", "sldMaster"),
            "notes": ("ppt/notesSlides/notesSlide1.xml", "ppt/notesSlides/_rels/notesSlide1.xml.rels", "notes"),
        }
        part, relation_part, root = owners[shared_image_owner]
        entries[part] = (
            f'<p:{root} xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}"><p:cSld><p:spTree>'
            '<p:pic><p:nvPicPr><p:cNvPr id="88" name="Shared image"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
            '<p:blipFill><a:blip r:embed="rIdSharedImage"/></p:blipFill></p:pic>'
            f'</p:spTree></p:cSld></p:{root}>'
        ).encode()
        existing = entries.get(relation_part, _rels())
        entries[relation_part] = existing.replace(
            b"</Relationships>",
            _rel("rIdSharedImage", "image", "../media/image1.png").encode() + b"</Relationships>",
        )
    with zipfile.ZipFile(path, "w") as archive:
        for name, value in entries.items():
            info = zipfile.ZipInfo(name, (2025, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            archive.writestr(info, value)
    return path


def _hashes(path: Path) -> dict[str, str]:
    with zipfile.ZipFile(path) as archive:
        return {name: hashlib.sha256(archive.read(name)).hexdigest() for name in archive.namelist()}


def _only_changed(source: Path, output: Path, *changed: str) -> None:
    before, after = _hashes(source), _hashes(output)
    assert before.keys() == after.keys()
    assert {name for name in before if before[name] != after[name]} == set(changed)


def test_inventory_exposes_stable_typed_locators_and_all_complex_parts(tmp_path: Path) -> None:
    source = _package(tmp_path / "real.pptx")
    assert len(Presentation(str(source)).slides) == 2
    inventory = PptxAdapter().inventory(source)
    assert [item["part_uri"] for item in inventory["slides"]] == ["ppt/slides/slide1.xml", "ppt/slides/slide2.xml"]
    assert inventory["placeholders"] == [{"part_uri": "ppt/slides/slide1.xml", "shape_id": "2", "placeholder_idx": "7"}]
    assert inventory["tables"][0]["shape_id"] == "3"
    assert inventory["charts"][0]["shape_id"] == "5"
    assert inventory["images"][0]["target_part"] == "ppt/media/image1.png"
    assert inventory["slide_size"] == {"width_emu": 9_144_000, "height_emu": 6_858_000}


def test_text_table_notes_and_image_edits_are_bounded_and_preserve_formatting(tmp_path: Path) -> None:
    source = _package(tmp_path / "source.pptx")
    adapter = PptxAdapter()
    text = tmp_path / "text.pptx"
    evidence = adapter.patch_text(source, text, {"part_uri": "ppt/slides/slide1.xml", "shape_id": "2"}, "New title", expected_text="BRAND TITLE")
    _only_changed(source, text, "ppt/slides/slide1.xml")
    with zipfile.ZipFile(text) as archive:
        xml = archive.read("ppt/slides/slide1.xml")
    assert b"New title" in xml and b'<a:rPr b="1"/>' in xml and b'<a:rPr i="1"/>' in xml
    preservation = cast(dict[str, object], evidence["preservation"])
    visual = cast(dict[str, object], evidence["visual_verification"])
    assert preservation["relationships_preserved"] is True
    assert visual["state"] == "not_run"

    table = tmp_path / "table.pptx"
    _ = adapter.patch_table_cell(source, table, {"part_uri": "ppt/slides/slide1.xml", "shape_id": "3", "row_index": 0, "column_index": 1}, "Changed", expected_text="B1")
    _only_changed(source, table, "ppt/slides/slide1.xml")

    notes = tmp_path / "notes.pptx"
    _ = adapter.patch_notes(source, notes, {"part_uri": "ppt/notesSlides/notesSlide1.xml", "shape_id": "9"}, "Private note", expected_text="Speaker note")
    _only_changed(source, notes, "ppt/notesSlides/notesSlide1.xml")

    image = tmp_path / "image.pptx"
    image_evidence = adapter.replace_image(source, image, {"part_uri": "ppt/slides/slide1.xml", "shape_id": "4"}, b"\x89PNG\r\n\x1a\nNEW")
    _only_changed(source, image, "ppt/media/image1.png")
    assert image_evidence["preservation"]["media_preserved"] is False


def test_aliased_duplicate_presentationml_shape_id_is_ambiguous(
    tmp_path: Path,
) -> None:
    source = _package(tmp_path / "duplicate-shape.pptx")
    with zipfile.ZipFile(source) as archive:
        parts = {name: archive.read(name) for name in archive.namelist()}
    namespace = b"http://schemas.openxmlformats.org/presentationml/2006/ma&#x69;n"
    duplicate = (
        b'<review:sp xmlns:review="' + namespace + b'">'
        b'<review:nvSpPr><review:cNvPr id="2" name="Duplicate title"/>'
        b'<review:cNvSpPr/><review:nvPr/></review:nvSpPr>'
        b'<review:txBody><a:p><a:r><a:t>DUPLICATE</a:t></a:r></a:p>'
        b'</review:txBody></review:sp>'
    )
    parts["ppt/slides/slide1.xml"] = parts["ppt/slides/slide1.xml"].replace(
        b'<p:transition advClick="1"/>',
        duplicate + b'<p:transition advClick="1"/>',
    )
    source = _zip(tmp_path / "semantic-duplicate-shape.pptx", parts)
    _ = Presentation(str(source))
    before = source.read_bytes()
    output = tmp_path / "ambiguous.pptx"

    with pytest.raises(DocumentError) as caught:
        _ = PptxAdapter().patch_text(
            source,
            output,
            {"part_uri": "ppt/slides/slide1.xml", "shape_id": "2"},
            "unsafe",
            expected_text="BRAND TITLE",
        )

    assert caught.value.code is DocumentErrorCode.AMBIGUOUS_LOCATOR
    assert caught.value.stage == "locate"
    assert source.read_bytes() == before
    assert not output.exists()
    _ = Presentation(str(source))


def test_strict_presentationml_and_drawingml_shapes_and_tables_are_targetable(
    tmp_path: Path,
) -> None:
    source = _package(tmp_path / "transitional.pptx")
    with zipfile.ZipFile(source) as archive:
        parts = {name: archive.read(name) for name in archive.namelist()}
    parts["ppt/slides/slide1.xml"] = (
        parts["ppt/slides/slide1.xml"]
        .replace(P.encode(), STRICT_P.encode())
        .replace(A.encode(), STRICT_A.encode())
    )
    source = _zip(tmp_path / "strict.pptx", parts)
    before = source.read_bytes()
    inventory = PptxAdapter().inventory(source)
    assert {item["shape_id"] for item in inventory["shapes"]} >= {"2", "3"}
    assert {item["shape_id"] for item in inventory["tables"]} == {"3"}

    text_output = tmp_path / "strict-text.pptx"
    _ = PptxAdapter().patch_text(
        source,
        text_output,
        {"part_uri": "ppt/slides/slide1.xml", "shape_id": "2"},
        "Strict title",
        expected_text="BRAND TITLE",
    )
    table_output = tmp_path / "strict-table.pptx"
    _ = PptxAdapter().patch_table_cell(
        source,
        table_output,
        {
            "part_uri": "ppt/slides/slide1.xml",
            "shape_id": "3",
            "row_index": 0,
            "column_index": 1,
        },
        "Strict cell",
        expected_text="B1",
    )

    with zipfile.ZipFile(text_output) as archive:
        assert b"Strict title" in archive.read("ppt/slides/slide1.xml")
    with zipfile.ZipFile(table_output) as archive:
        assert b"Strict cell" in archive.read("ppt/slides/slide1.xml")
    assert source.read_bytes() == before


def test_generic_shape_text_patch_refuses_table_graphic_frame(tmp_path: Path) -> None:
    source = _package(tmp_path / "table-source.pptx")
    output = tmp_path / "refused.pptx"

    with pytest.raises(DocumentError) as caught:
        _ = PptxAdapter().patch_text(
            source,
            output,
            {"part_uri": "ppt/slides/slide1.xml", "shape_id": "3"},
            "unsafe whole-table rewrite",
        )

    assert caught.value.code is DocumentErrorCode.UNSUPPORTED_EDIT
    assert caught.value.details["reason"] == "table_shape_requires_cell_locator"
    assert not output.exists()


def test_generic_shape_text_patch_refuses_alias_prefixed_table(tmp_path: Path) -> None:
    source = _package(tmp_path / "aliased-table.pptx")
    with zipfile.ZipFile(source) as archive:
        parts = {name: archive.read(name) for name in archive.namelist()}
    slide = parts["ppt/slides/slide1.xml"]
    slide = slide.replace(
        b"<a:graphicData>",
        f'<a:graphicData xmlns:d="{A}" uri="http://schemas.openxmlformats.org/drawingml/2006/table">'.encode(),
        1,
    ).replace(b"<a:tbl>", b"<d:tbl>", 1).replace(b"</a:tbl>", b"</d:tbl>", 1)
    parts["ppt/slides/slide1.xml"] = slide
    source = _zip(tmp_path / "semantic-table.pptx", parts)
    reopened = Presentation(str(source))
    assert next(shape for shape in reopened.slides[0].shapes if shape.shape_id == 3).has_table

    output = tmp_path / "refused.pptx"
    with pytest.raises(DocumentError) as caught:
        _ = PptxAdapter().patch_text(
            source,
            output,
            {"part_uri": "ppt/slides/slide1.xml", "shape_id": "3"},
            "unsafe whole-table rewrite",
        )

    assert caught.value.code is DocumentErrorCode.UNSUPPORTED_EDIT
    assert caught.value.details["reason"] == "table_shape_requires_cell_locator"
    assert not output.exists()


def test_image_replacement_refuses_two_shapes_consuming_one_relationship(
    tmp_path: Path,
) -> None:
    source = _package(tmp_path / "one-relationship.pptx")
    with zipfile.ZipFile(source) as archive:
        parts = {name: archive.read(name) for name in archive.namelist()}
    second_picture = (
        b'<p:pic><p:nvPicPr><p:cNvPr id="6" name="Second image"/>'
        b'<p:cNvPicPr/><p:nvPr/></p:nvPicPr><p:blipFill>'
        b'<a:blip r:embed="rIdImage"/></p:blipFill></p:pic>'
    )
    parts["ppt/slides/slide1.xml"] = parts["ppt/slides/slide1.xml"].replace(
        b'<p:transition advClick="1"/>',
        second_picture + b'<p:transition advClick="1"/>',
    )
    source = _zip(tmp_path / "shared-rid.pptx", parts)
    reopened = Presentation(str(source))
    assert len(reopened.slides[0].shapes) == 5
    consumers = [
        item for item in PptxAdapter().inventory(source)["images"]
        if item["target_part"] == "ppt/media/image1.png"
    ]
    assert {item["shape_id"] for item in consumers} == {"4", "6"}

    output = tmp_path / "refused.pptx"
    with pytest.raises(DocumentError) as caught:
        _ = PptxAdapter().replace_image(
            source,
            output,
            {"part_uri": "ppt/slides/slide1.xml", "shape_id": "4"},
            b"\x89PNG\r\n\x1a\nNEW",
        )

    assert caught.value.code is DocumentErrorCode.LOSSY_WRITE_BLOCKED
    assert caught.value.details == {"reason": "shared_media_target", "references": 2}
    assert not output.exists()


@pytest.mark.parametrize("owner", ["layout", "master", "notes"])
def test_image_replacement_refuses_target_shared_outside_selected_slide(
    tmp_path: Path, owner: str
) -> None:
    source = _package(tmp_path / f"shared-{owner}.pptx", shared_image_owner=owner)
    inventory = PptxAdapter().inventory(source)
    references = [
        item for item in inventory["images"]
        if item["target_part"] == "ppt/media/image1.png"
    ]
    assert {item["part_uri"] for item in references} == {
        "ppt/slides/slide1.xml",
        {
            "layout": "ppt/slideLayouts/slideLayout1.xml",
            "master": "ppt/slideMasters/slideMaster1.xml",
            "notes": "ppt/notesSlides/notesSlide1.xml",
        }[owner],
    }

    output = tmp_path / "refused.pptx"
    with pytest.raises(DocumentError) as caught:
        _ = PptxAdapter().replace_image(
            source,
            output,
            {"part_uri": "ppt/slides/slide1.xml", "shape_id": "4"},
            b"\x89PNG\r\n\x1a\nNEW",
        )
    assert caught.value.code is DocumentErrorCode.LOSSY_WRITE_BLOCKED
    assert caught.value.details == {"reason": "shared_media_target", "references": 2}
    assert not output.exists()


@pytest.mark.parametrize("consumer", ["background", "extension"])
def test_image_replacement_counts_every_xml_embed_consumer(
    tmp_path: Path, consumer: str
) -> None:
    source = _package(tmp_path / f"{consumer}-consumer.pptx")
    with zipfile.ZipFile(source) as archive:
        parts = {name: archive.read(name) for name in archive.namelist()}
    slide = parts["ppt/slides/slide1.xml"]
    if consumer == "background":
        extra = (
            b'<p:bg><p:bgPr><a:blipFill><a:blip r:embed="rIdImage"/>'
            b"<a:stretch><a:fillRect/></a:stretch></a:blipFill>"
            b"<a:effectLst/></p:bgPr></p:bg>"
        )
        slide = slide.replace(b"<p:spTree>", extra + b"<p:spTree>", 1)
    else:
        extra = (
            b'<p:extLst><p:ext uri="{BIRKIN-EMBED-CONSUMER}">'
            b'<p14:media xmlns:p14="http://schemas.microsoft.com/office/'
            b'powerpoint/2010/main" r:embed="rIdImage"/>'
            b"</p:ext></p:extLst>"
        )
        slide = slide.replace(b"</p:sld>", extra + b"</p:sld>", 1)
    parts["ppt/slides/slide1.xml"] = slide
    source = _zip(tmp_path / f"semantic-{consumer}-consumer.pptx", parts)
    _ = Presentation(str(source))
    before = source.read_bytes()
    output = tmp_path / "refused.pptx"

    with pytest.raises(DocumentError) as caught:
        _ = PptxAdapter().replace_image(
            source,
            output,
            {"part_uri": "ppt/slides/slide1.xml", "shape_id": "4"},
            b"\x89PNG\r\n\x1a\nNEW",
        )

    assert caught.value.code is DocumentErrorCode.LOSSY_WRITE_BLOCKED
    assert caught.value.details == {
        "reason": "shared_media_target",
        "references": 2,
    }
    assert source.read_bytes() == before
    assert not output.exists()
    _ = Presentation(str(source))


def test_reorder_and_page_size_preserve_transitions_comments_and_graph(tmp_path: Path) -> None:
    source = _package(tmp_path / "source.pptx", sections=False)
    adapter = PptxAdapter()
    reordered = tmp_path / "reordered.pptx"
    evidence = adapter.reorder_slides(source, reordered, ["ppt/slides/slide2.xml", "ppt/slides/slide1.xml"])
    _only_changed(source, reordered, "ppt/presentation.xml")
    with zipfile.ZipFile(reordered) as archive:
        presentation = archive.read("ppt/presentation.xml")
        assert presentation.index(b'r:id="rId2"') < presentation.index(b'r:id="rId1"')
    assert evidence["preservation"]["relationships_preserved"] is True

    sized = tmp_path / "sized.pptx"
    _ = adapter.set_page_size(source, sized, 12_192_000, 6_858_000)
    _only_changed(source, sized, "ppt/presentation.xml")
    with zipfile.ZipFile(sized) as archive:
        presentation = archive.read("ppt/presentation.xml")
    assert b'cx="12192000"' in presentation and b'type="screen4x3"' in presentation


def test_reorder_refuses_section_boundary_rewrite(tmp_path: Path) -> None:
    source = _package(tmp_path / "sectioned.pptx")
    output = tmp_path / "refused.pptx"
    with pytest.raises(DocumentError) as caught:
        _ = PptxAdapter().reorder_slides(source, output, ["ppt/slides/slide2.xml", "ppt/slides/slide1.xml"])
    assert caught.value.code is DocumentErrorCode.LOSSY_WRITE_BLOCKED
    assert caught.value.details["reason"] == "section_boundary_rewrite_required"
    assert not output.exists()


@pytest.mark.parametrize("method", ["add_slide", "delete_slide", "change_layout", "update_chart_data", "replace_linked_media", "update_theme", "update_master"])
def test_graph_wide_operations_refuse_without_output(method: str, tmp_path: Path) -> None:
    source = _package(tmp_path / "source.pptx")
    output = tmp_path / "refused.pptx"
    with pytest.raises(DocumentError) as caught:
        getattr(PptxAdapter(), method)(source, output)
    assert caught.value.code is DocumentErrorCode.LOSSY_WRITE_BLOCKED
    assert caught.value.details["status"] == "refused"
    assert not output.exists()
