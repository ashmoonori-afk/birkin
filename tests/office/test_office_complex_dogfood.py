from __future__ import annotations

import hashlib
import json
import shutil
import subprocess
import sys
from pathlib import Path
from typing import cast
from zipfile import ZipFile

from defusedxml import ElementTree
from docx import Document
from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pypdf import PdfReader

ROOT = Path(__file__).parents[2]
SCRIPT = ROOT / "script/qa/office_complex_dogfood.py"
FORMATS = {"docx", "xlsx", "pptx", "pdf", "hwpx"}
TOOLS = {
    "apply_document_patch",
    "compare_documents",
    "convert_document",
    "create_document",
    "extract_document",
    "fill_template",
    "inspect_document",
    "list_document_adapters",
    "render_artifact",
    "validate_artifact",
}


def _invoke(output: Path, *, module: bool) -> dict[str, object]:
    command = [sys.executable]
    command += ["-m", "script.qa.office_complex_dogfood"] if module else [str(SCRIPT)]
    command += ["--output-dir", str(output)]
    completed = subprocess.run(command, cwd=ROOT, text=True, capture_output=True, check=False)
    assert completed.returncode == 0, completed.stderr or completed.stdout
    assert len(completed.stdout.encode()) < 20_000
    summary = cast(dict[str, object], json.loads(completed.stdout))
    assert summary["ok"] is True
    assert Path(cast(str, summary["jail"])).resolve() == output.resolve()
    summary_formats = cast(dict[str, dict[str, object]], summary["formats"])
    assert all(item["source_preserved"] is True for item in summary_formats.values())

    full = cast(dict[str, object], summary["full_report"])
    report_path = Path(cast(str, full["uri"]))
    payload = report_path.read_bytes()
    assert len(payload) == full["bytes"]
    assert hashlib.sha256(payload).hexdigest() == full["sha256"]
    report = cast(dict[str, object], json.loads(payload))
    assert report["ok"] is True
    assert set(cast(list[str], report["tools_exercised"])) == TOOLS
    return report


def test_literal_complex_office_dogfood_path_and_module(tmp_path: Path) -> None:
    path_output = tmp_path / "path-evidence"
    module_output = tmp_path / "module-evidence"
    report = _invoke(path_output, module=False)
    second = _invoke(module_output, module=True)
    assert cast(dict[str, object], second["formats"]).keys() == cast(dict[str, object], report["formats"]).keys()

    formats = cast(dict[str, dict[str, object]], report["formats"])
    assert set(formats) == FORMATS
    for evidence in formats.values():
        assert evidence["source_sha256_before"] == evidence["source_sha256_after"]
        operations = cast(dict[str, object], evidence["operations"])
        assert {"create", "inspect", "extract", "modify", "validate", "diff", "structured_preview", "convert_txt", "visual_render"} == set(operations)
        preview = cast(dict[str, object], operations["structured_preview"])
        visual = cast(dict[str, object], operations["visual_render"])
        assert preview == {"status": "preview", "visual_proof": False}
        assert visual["status"] == "unavailable"
        artifact = Path(cast(str, evidence["modified_artifact"]))
        assert artifact.is_file() and artifact.resolve().is_relative_to(path_output.resolve())
        reopen = cast(dict[str, object], evidence["reopen_validation"])
        assert reopen["status"] == "ok"
        assert evidence["expected_content_match"] is True

    docx_path = Path(cast(str, formats["docx"]["modified_artifact"]))
    document = Document(str(docx_path))
    assert len(document.tables) == 1
    with ZipFile(docx_path) as package:
        document_xml = package.read("word/document.xml")
        assert document_xml.count(b'w:type="page"') >= 2
        assert b"commentRangeStart" in document_xml and b"<w:ins" in document_xml
        assert b"fldSimple" in document_xml and b"Modified customer" in document_xml
        assert package.read("word/comments.xml") and package.read("word/header1.xml")

    xlsx_path = Path(cast(str, formats["xlsx"]["modified_artifact"]))
    workbook = load_workbook(xlsx_path, data_only=False, read_only=False)
    data_sheet = cast(Worksheet, workbook["Data"])
    assert cast(object, data_sheet["B2"].value) == 42
    assert workbook["Hidden"].sheet_state == "hidden"
    assert workbook["VeryHidden"].sheet_state == "veryHidden"
    assert data_sheet.row_dimensions[3].hidden
    assert data_sheet.column_dimensions["D"].hidden
    workbook.close()
    with ZipFile(xlsx_path) as package:
        assert any(name.startswith("xl/charts/") for name in package.namelist())
    formula = cast(dict[str, object], formats["xlsx"]["formula_evidence"])
    assert cast(int, formula["count"]) >= 4
    assert cast(list[object], formula["dynamic"])
    assert cast(dict[str, object], formula["external"])
    assert cast(dict[str, object], formula["recalculation"])["state"] == "unavailable"

    pptx_path = Path(cast(str, formats["pptx"]["modified_artifact"]))
    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) == 2
    assert presentation.slides[0].notes_slide.notes_text_frame is not None
    with ZipFile(pptx_path) as package:
        names = set(package.namelist())
        assert any(name.startswith("ppt/slideMasters/") for name in names)
        assert any(name.startswith("ppt/slideLayouts/") for name in names)
        assert any(name.startswith("ppt/theme/") for name in names)
        assert any(name.startswith("ppt/notesSlides/") for name in names)
        assert any(name.startswith("ppt/media/") for name in names)
    layout = cast(dict[str, object], formats["pptx"]["layout_evidence"])
    assert cast(list[object], layout["warnings"])
    assert cast(dict[str, object], layout["visual_verification"])["status"] == "not-run"

    pdf_path = Path(cast(str, formats["pdf"]["modified_artifact"]))
    assert len(PdfReader(pdf_path, strict=True).pages) == 1
    pdf_states = cast(dict[str, dict[str, object]], formats["pdf"]["state_cases"])
    assert formats["pdf"]["inspection_surface"] == "registered_inspect_document"
    assert set(pdf_states) == {"native", "scanned", "form", "active", "signed", "encrypted"}
    assert pdf_states["scanned"]["content_type"] == "image_only"
    assert pdf_states["form"]["form_type"] == "acroform"
    assert pdf_states["active"]["active_content"]
    assert pdf_states["signed"]["signed"] is True
    assert pdf_states["encrypted"]["credential_required"] is True
    refusal_reasons: set[str] = set()
    for raw_refusal in cast(list[object], formats["pdf"]["typed_refusals"]):
        refusal = cast(dict[str, object], raw_refusal)
        details = cast(dict[str, object], refusal["details"])
        refusal_reasons.add(str(details["reason"]))
    assert {"pdf_image_only_requires_ocr", "pdf_acroform_fill_unsupported", "pdf_password_required"} <= refusal_reasons

    hwpx_path = Path(cast(str, formats["hwpx"]["modified_artifact"]))
    with ZipFile(hwpx_path) as package:
        assert package.read("mimetype") == b"application/hwp+zip"
        sections = [name for name in package.namelist() if name.startswith("Contents/section")]
        assert len(sections) == 2
        section = package.read("Contents/section0.xml")
        _ = ElementTree.fromstring(section)
        assert b"<hp:tbl" in section and b"Modified customer" in section

    unsupported = cast(dict[str, dict[str, object]], report["unsupported_identities"])
    assert set(unsupported) == {"odt", "ods", "odp", "hwp"}
    assert all(item["status"] == "unsupported" for item in unsupported.values())
    for ext in ("odt", "ods", "odp"):
        with ZipFile(Path(cast(str, unsupported[ext]["path"]))) as package:
            assert package.read("mimetype").startswith(b"application/vnd.oasis.opendocument")
            _ = ElementTree.fromstring(package.read("content.xml"))
    hwp_path = Path(cast(str, unsupported["hwp"]["path"]))
    assert hwp_path.read_bytes().startswith(bytes.fromhex("d0cf11e0a1b11ae1"))
    strict_hwp = cast(dict[str, object], unsupported["hwp"]["strict_identity_refusal"])
    assert strict_hwp["status"] == "refused" and strict_hwp["reason_code"] == "invalid_cfb_identity"
    legacy = cast(dict[str, dict[str, object]], report["legacy"])
    assert set(legacy) == {"doc", "xls", "ppt", "rtf"}
    for item in legacy.values():
        identity = cast(dict[str, object], item["identity"])
        conversion = cast(dict[str, object], item["conversion"])
        assert identity["status"] == "accepted"
        assert conversion["status"] == "refused"
        assert conversion["reason_code"] == "external_engine_forbidden"
        assert item["source_immutable"] is True

    cleanup = cast(dict[str, object], report["cleanup_receipt"])
    assert cleanup["temporary_paths_remaining"] == []
    assert cleanup["artifacts_retained"] is True
    for removed in cast(list[str], cleanup["temporary_paths_removed"]):
        assert not Path(removed).exists()
    shutil.rmtree(path_output)
    shutil.rmtree(module_output)
    assert not path_output.exists() and not module_output.exists()
