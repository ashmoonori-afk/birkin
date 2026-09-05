from __future__ import annotations

from pathlib import Path
from typing import cast

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from birkin.office.errors import DocumentError
from birkin.office.service import DocumentService


def _apply(
    service: DocumentService,
    artifact: dict[str, str],
    operations: list[dict[str, object]],
    output_name: str,
) -> Path:
    result = service.apply_document_patch(
        artifact,
        {"operations": operations},
        expected_source_sha256=artifact["content_hash"],
        output_name=output_name,
        dry_run=False,
    )
    return Path(cast("dict[str, str]", result["draft_artifact"])["uri"])


def test_multiple_docx_paragraphs_publish_as_one_draft(tmp_path: Path) -> None:
    service = DocumentService(tmp_path)
    created = service.create_document(
        format="docx", content={"paragraphs": ["one", "two", "three"]}, output_name="source.docx"
    )
    artifact = cast("dict[str, str]", created["draft_artifact"])

    output = _apply(service, artifact, [
        {"locator": {"format": "docx", "index": 1}, "value": "ONE"},
        {"locator": {"format": "docx", "index": 3}, "value": "THREE"},
    ], "patched.docx")

    assert [paragraph.text for paragraph in Document(str(output)).paragraphs] == ["ONE", "two", "THREE"]


def test_multiple_xlsx_sheets_use_named_locators_and_fail_atomically(tmp_path: Path) -> None:
    service = DocumentService(tmp_path)
    created = service.create_document(format="xlsx", content={"sheets": [
        {"name": "January", "rows": [["Value"], [1]]},
        {"name": "February", "rows": [["Value"], [2]]},
    ]}, output_name="source.xlsx")
    artifact = cast("dict[str, str]", created["draft_artifact"])
    source = Path(artifact["uri"])
    before = source.read_bytes()

    output = _apply(service, artifact, [
        {"locator": {"sheet": "January", "cell": "A2"}, "value": 10},
        {"locator": {"sheet": "February", "cell": "A2"}, "value": 20},
    ], "patched.xlsx")
    workbook = load_workbook(output, read_only=True)
    assert workbook["January"]["A2"].value == 10
    assert workbook["February"]["A2"].value == 20

    with pytest.raises(DocumentError):
        _apply(service, artifact, [
            {"locator": {"sheet": "January", "cell": "A2"}, "value": 30},
            {"locator": {"sheet": "Missing", "cell": "A2"}, "value": 40},
        ], "failed.xlsx")
    assert source.read_bytes() == before
    assert not (tmp_path / "artifacts" / "drafts" / "failed.xlsx").exists()


def test_multiple_pptx_slides_use_part_locators(tmp_path: Path) -> None:
    service = DocumentService(tmp_path)
    created = service.create_document(format="pptx", content={"slides": [
        {"title": "One", "body": "old one"},
        {"title": "Two", "body": "old two"},
    ]}, output_name="source.pptx")
    artifact = cast("dict[str, str]", created["draft_artifact"])

    output = _apply(service, artifact, [
        {"locator": {"slide_part": "ppt/slides/slide1.xml", "placeholder_idx": 1}, "value": "new one"},
        {"locator": {"slide_part": "ppt/slides/slide2.xml", "placeholder_idx": 1}, "value": "new two"},
    ], "patched.pptx")
    slides = Presentation(str(output)).slides
    assert slides[0].placeholders[1].text == "new one"
    assert slides[1].placeholders[1].text == "new two"
